#!/usr/bin/env python3
"""Extract document TEXT from the Drive into the private repo (Phase 7c).

Reads the manifest (`<private-repo>/drive/drive-index.json`), and for each
content-bearing file under a subtree prefix, pulls its **text** — Google Docs /
Slides via headless API export, Word/PowerPoint/PDF/plain via download + a local
extractor — into `<private-repo>/drive/extracts/<mirrored path>/<title>__<id8>.txt`
with a provenance header. INTERNAL, like the manifest (docs/PRIVACY.md D44b/D46).

Idempotent + resumable: an index (`drive/.extract-index.json`) records each file's
`modified`, so a re-run only (re-)extracts new/changed files. Per-file failures are
logged and skipped, never fatal. Scanned PDFs (no text layer) are recorded as empty
so they aren't retried.

Auth + venv: same as gen_drive_index.py (read-only Drive ADC; the `~/.config/ds-kb`
venv, which also needs `python-docx python-pptx`; `pdftotext` on PATH for PDFs).

Usage:
  GOOGLE_APPLICATION_CREDENTIALS= ~/.config/ds-kb/venv/bin/python scripts/gen_drive_extracts.py \
      [--prefix "CERF Anticipatory Action"] [--types doc,slides,document,presentation,plain] \
      [--limit N] [--max-pdf-mb 40]
  --types: comma list of friendly types, or `all`. Default = the narrative set
           (no pdf/sheet). Add `pdf` for the heavy PDF pass.
"""
from __future__ import annotations
import io, json, os, re, subprocess, sys, tempfile, threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
INTERNAL = Path(os.environ.get("KB_INTERNAL_DIR", ROOT.parent / "ds-knowledge-base-internal"))
MANIFEST = INTERNAL / "drive" / "drive-index.json"
OUTDIR = INTERNAL / "drive" / "extracts"
INDEX = INTERNAL / "drive" / ".extract-index.json"

NARRATIVE = ["doc", "slides", "document", "presentation", "plain"]
EXPORT_MIME = {"doc": "text/plain", "slides": "text/plain", "sheet": "text/csv"}  # Google-native
MIN_CHARS = 20            # below this an extract is treated as empty (e.g. scanned PDF)


def arg(name, default=None):
    return sys.argv[sys.argv.index(name) + 1] if name in sys.argv else default


_CREDS = None             # resolved once in main()
_tl = threading.local()   # one Drive service per worker thread (httplib2 isn't thread-safe)


def make_creds():
    scopes = ["https://www.googleapis.com/auth/drive.readonly"]
    cred = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS")
    if cred:
        from google.oauth2 import service_account
        return service_account.Credentials.from_service_account_file(cred, scopes=scopes)
    import google.auth
    creds, _ = google.auth.default(scopes=scopes)
    return creds


def svc():
    if not hasattr(_tl, "s"):
        from googleapiclient.discovery import build
        _tl.s = build("drive", "v3", credentials=_CREDS, cache_discovery=False)
    return _tl.s


def sanitize(s: str, maxlen=120) -> str:
    s = re.sub(r"[^\w\-. ]+", "_", (s or "").strip()).strip("._ ")
    return (s or "untitled")[:maxlen]


# ---- per-type text extractors -------------------------------------------------
def from_docx(data: bytes) -> str:
    import docx
    d = docx.Document(io.BytesIO(data))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def from_pptx(data: bytes) -> str:
    import pptx
    pres = pptx.Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(pres.slides, 1):
        parts.append(f"--- slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def from_pdf(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as f:
        f.write(data); f.flush()
        r = subprocess.run(["pdftotext", "-enc", "UTF-8", f.name, "-"],
                           capture_output=True, timeout=120)
    return r.stdout.decode("utf-8", "ignore")


def extract_text(node) -> str:
    t = node["type"]
    fid = node["id"]
    if t in EXPORT_MIME:                       # Google-native → export
        data = svc().files().export_media(fileId=fid, mimeType=EXPORT_MIME[t]).execute()
        return data.decode("utf-8", "ignore") if isinstance(data, bytes) else str(data)
    data = svc().files().get_media(fileId=fid).execute()   # binary → download
    if t == "document":
        return from_docx(data)
    if t == "presentation":
        return from_pptx(data)
    if t == "pdf":
        return from_pdf(data)
    if t == "plain":
        return data.decode("utf-8", "ignore")
    raise ValueError(f"no extractor for type {t}")


def worker(n):
    """Extract + write one file (runs in a thread). Returns a result dict; never raises."""
    fid, t, modn = n["id"], n["type"], n.get("modified", "")
    try:
        text = extract_text(n).strip()
    except Exception as e:                     # noqa: BLE001 — log + continue
        return {"fid": fid, "status": "fail",
                "msg": f"[{t}] {n.get('path','')}/{n['title']}: {type(e).__name__}: {str(e)[:80]}"}
    if len(text) < MIN_CHARS:
        return {"fid": fid, "status": "empty", "record": {"modified": modn, "empty": True, "reason": "no-text"}}
    rel = Path(sanitize_path(n.get("path", ""))) / f"{sanitize(n['title'])}__{fid[:8]}.txt"
    out = OUTDIR / rel
    out.parent.mkdir(parents=True, exist_ok=True)
    header = (f"# {n['title']}\n# drive-path: {n.get('path','')}\n"
              f"# type: {t} · modified: {modn} · id: {fid}\n# link: {n.get('link','')}\n"
              f"{'-'*60}\n")
    out.write_text(header + text)
    return {"fid": fid, "status": "done", "record": {"modified": modn, "out": str(rel)}}


def main():
    global _CREDS
    if not MANIFEST.exists():
        sys.exit(f"Manifest not found at {MANIFEST} — run gen_drive_index.py first.")
    prefix = arg("--prefix", "CERF Anticipatory Action")
    types_arg = arg("--types", ",".join(NARRATIVE))
    types = (["doc", "slides", "document", "presentation", "plain", "pdf", "sheet"]
             if types_arg == "all" else [x.strip() for x in types_arg.split(",")])
    limit = int(arg("--limit", "0"))
    max_pdf_mb = float(arg("--max-pdf-mb", "40"))
    workers = int(arg("--workers", "10"))

    nodes = json.loads(MANIFEST.read_text())["nodes"]
    sel = [n for n in nodes if n["type"] in types and not n.get("excluded")
           and (n.get("path", "") + "/").startswith(prefix)]
    if limit:
        sel = sel[:limit]
    index = json.loads(INDEX.read_text()) if INDEX.exists() else {}

    def size_mb(s):
        m = re.match(r"([\d.]+)(KB|MB|GB|B)", s or "")
        if not m: return 0.0
        return float(m.group(1)) * {"B": 1/1024/1024, "KB": 1/1024, "MB": 1, "GB": 1024}[m.group(2)]

    # pre-filter: drop already-extracted (idempotent) and oversize PDFs
    skip = 0
    todo = []
    for n in sel:
        fid, t, modn = n["id"], n["type"], n.get("modified", "")
        rec = index.get(fid)
        if rec and rec.get("modified") == modn and (rec.get("empty") or (OUTDIR / rec.get("out", "")).exists()):
            skip += 1; continue
        if t == "pdf" and size_mb(n.get("size", "")) > max_pdf_mb:
            index[fid] = {"modified": modn, "empty": True, "reason": "pdf>max"}; skip += 1; continue
        todo.append(n)

    _CREDS = make_creds()
    done = empty = fail = 0
    total = len(todo)
    print(f"{len(sel)} candidates under '{prefix}' (types: {','.join(types)}); "
          f"{skip} already done/skipped, {total} to extract with {workers} workers.")
    with ThreadPoolExecutor(max_workers=workers) as ex:
        futs = [ex.submit(worker, n) for n in todo]
        for i, fut in enumerate(as_completed(futs), 1):
            r = fut.result()
            if r["status"] == "fail":
                fail += 1
                print(f"  FAIL {r['msg']}")
            else:
                index[r["fid"]] = r["record"]
                done += r["status"] == "done"
                empty += r["status"] == "empty"
            if i % 100 == 0:
                INDEX.write_text(json.dumps(index, indent=0, ensure_ascii=False))
                print(f"  …{i}/{total}  (extracted {done}, empty {empty}, fail {fail})")
    INDEX.write_text(json.dumps(index, indent=0, ensure_ascii=False))
    print(f"done: {done} extracted / {empty} empty (no text) / {skip} skipped / {fail} failed → {OUTDIR}")


def sanitize_path(path: str) -> str:
    return "/".join(sanitize(seg) for seg in path.split("/") if seg)


if __name__ == "__main__":
    main()
